import os
import unicodedata
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
import PyPDF2
import openpyxl
import PySimpleGUI as sg
import win32com.client
import pandas as pd

file_paths = {}

def normalize_text(text):
    normalized_text = unicodedata.normalize('NFKC', text)
    normalized_text = normalized_text.replace('ｶﾞ', 'ガ').replace('ｷﾞ', 'ギ').replace('ｸﾞ', 'グ').replace('ｹﾞ', 'ゲ').replace('ｺﾞ', 'ゴ')
    # 他の半角カナの変換ルールも同様に追加
    normalized_text = normalized_text.replace('ｱ', 'ア').replace('ｲ', 'イ').replace('ｳ', 'ウ').replace('ｴ', 'エ').replace('ｵ', 'オ')
    normalized_text = normalized_text.replace('ｻﾞ', 'ザ').replace('ｼﾞ', 'ジ').replace('ｽﾞ', 'ズ').replace('ｾﾞ', 'ゼ').replace('ｿﾞ', 'ゾ')
    normalized_text = normalized_text.replace('ﾀﾞ', 'ダ').replace('ﾁﾞ', 'ヂ').replace('ﾂﾞ', 'ヅ').replace('ﾃﾞ', 'デ').replace('ﾄﾞ', 'ド')
    normalized_text = normalized_text.replace('ﾊﾞ', 'バ').replace('ﾋﾞ', 'ビ').replace('ﾌﾞ', 'ブ').replace('ﾍﾞ', 'ベ').replace('ﾎﾞ', 'ボ')
    normalized_text = normalized_text.replace('ﾊﾟ', 'パ').replace('ﾋﾟ', 'ピ').replace('ﾌﾟ', 'プ').replace('ﾍﾟ', 'ペ').replace('ﾎﾟ', 'ポ')
    normalized_text = normalized_text.replace('ｳﾞ', 'ヴ').replace('ﾜﾞ', 'ヷ').replace('ｦﾞ', 'ヺ').replace('ﾞ', '゛').replace('ﾟ', '゜')
    return normalized_text.lower()

def search_files(folder_path, keywords, exclude_keywords, selected_extensions=None):
    found_files = []
    total_files = sum([len(files) for _, _, files in os.walk(folder_path)])

    layout = [
        [sg.Text('検索中...')],
        [sg.ProgressBar(total_files, orientation='h', size=(20, 20), key='-PROGRESS-', bar_color=('blue', 'black'))],
    ]

    window = sg.Window('進捗状況', layout, finalize=True)
    progress_bar = window['-PROGRESS-']

    progress_bar.update_bar(0)
    current_file_index = 0
    
    for root, dirs, files in os.walk(folder_path):
        for file in files:
                # ファイルの拡張子が指定されたものであるかを確認
                if file.lower().endswith(('.pptx', '.docx', '.xlsx', '.pdf', '.doc', '.xls')):
                    # 選択された拡張子以外は無視
                    if selected_extensions and not any(file.lower().endswith(ext) for ext in selected_extensions):
                        continue
                    # ファイルのフルパスを取得
                    full_path = os.path.join(root, file)
                    # ファイル名を小文字に変換し、半角カナを全角カナに変換
                    file_name_normalized = normalize_text(file.lower())
                    # キーワードも半角カナを全角カナに変換
                    keywords_normalized = [normalize_text(kw) for kw in keywords]
                    # 除外キーワードも半角カナを全角カナに変換
                    exclude_keywords_normalized = [normalize_text(kw) for kw in exclude_keywords]

                    # ファイル名のみでキーワード検索
                    count_file_name = sum(file_name_normalized.count(kw) for kw in keywords_normalized)
                    if any(exclude_keywords_normalized):
                        if any(exclude_keyword and normalize_text(exclude_keyword.lower()) in file_name_normalized for exclude_keyword in exclude_keywords_normalized):
                            # 除外キーワードがある場合は候補から外す
                            continue
                    # キーワードがある場合はカウント＋１
                    count = count_file_name

                    # 拡張子ごとの関数でファイルを検索
                    if file_name_normalized.endswith('.pptx') and not file_name_normalized.startswith('~$'):
                        count_text, found = search_pptx(full_path, keywords_normalized, exclude_keywords_normalized)
                    elif file_name_normalized.endswith('.docx') and not file_name_normalized.startswith('~$'):
                        count_text, found = search_docx(full_path, keywords_normalized, exclude_keywords_normalized)
                    elif file.lower().endswith('.xlsx') and not file_name_normalized.startswith('~$'):
                        count_text, found = search_xlsx(full_path, keywords_normalized, exclude_keywords_normalized)
                    elif file.lower().endswith('.xls') and not file_name_normalized.startswith('~$'):
                        count_text, found = search_xls(full_path, keywords_normalized, exclude_keywords_normalized)
                    elif file_name_normalized.endswith('.pdf') and not file_name_normalized.startswith('~$'):
                        count_text, found = search_pdf(full_path, keywords_normalized, exclude_keywords_normalized)
                    elif file_name_normalized.endswith('.doc') and not file_name_normalized.startswith('~$'): 
                        count_text, found = search_doc(full_path, keywords_normalized, exclude_keywords_normalized)
                    else:
                        count_text, found = 0, any(
                            exclude_keyword and normalize_text(exclude_keyword.lower()) in file_name_normalized for exclude_keyword in
                            exclude_keywords_normalized)

                    # キーワードの数と検索結果を保存
                    count += count_text
                    if found or (count_file_name > 0 and not any(
                            exclude_keyword and normalize_text(exclude_keyword.lower()) in file_name_normalized for exclude_keyword in
                            exclude_keywords_normalized)):
                        found_files.append((full_path, count))

                        current_file_index += 1
                        progress_bar.UpdateBar(current_file_index)
    window.close()
    return found_files

def search_pptx(file_path, keywords, exclude_keywords):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
    text = '\n'.join(text_runs)
    text = normalize_text(text.lower())  # ファイルの中身にある半角カナを全角カナに変換
    keywords_normalized = [normalize_text(kw) for kw in keywords]
    exclude_keywords_normalized = [normalize_text(kw) for kw in exclude_keywords]
    
    # 除外キーワードがある場合はカウントを-99999にする
    count = 0
    if any(exclude_keywords_normalized):
        for exclude_keyword in exclude_keywords_normalized:
            if exclude_keyword and normalize_text(exclude_keyword.lower()) in text:
                return (count - 99999, False)

    count = sum(text.count(kw) for kw in keywords_normalized)
    
    # カウントが1以上の場合のみ検索結果に表示する
    return (count, any(text.find(kw) != -1 for kw in keywords_normalized))

def search_xls(file_path, keywords, exclude_keywords):
    text = ''
    
    try:
        df = pd.read_excel(file_path, header=None)
        text = df.to_string(header=False, index=False)
    except Exception as e:
        print(f"Error reading Excel file: {e}")

    text = normalize_text(text.lower())  # ファイルの中身にある半角カナを全角カナに変換
    keywords_normalized = [normalize_text(kw) for kw in keywords]
    exclude_keywords_normalized = [normalize_text(kw) for kw in exclude_keywords]

    # 除外キーワードがある場合はカウントを-99999にする
    count = 0
    if any(exclude_keywords_normalized):
        for exclude_keyword in exclude_keywords_normalized:
            if exclude_keyword and normalize_text(exclude_keyword.lower()) in text:
                return (count - 99999, False)

    count = sum(text.count(kw) for kw in keywords_normalized)

    # カウントが1以上の場合のみ検索結果に表示する
    return (count, any(text.find(kw) != -1 for kw in keywords_normalized))

def search_docx(file_path, keywords, exclude_keywords):
    doc = Document(file_path)
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text
    text = normalize_text(text.lower())  # ファイルの中身にある半角カナを全角カナに変換
    keywords_normalized = [normalize_text(kw) for kw in keywords]
    exclude_keywords_normalized = [normalize_text(kw) for kw in exclude_keywords]

    # 除外キーワードがある場合はカウントを-99999にする
    count = 0
    if any(exclude_keywords_normalized):
        for exclude_keyword in exclude_keywords_normalized:
            if exclude_keyword and normalize_text(exclude_keyword.lower()) in text:
                return (count - 99999, False)

    count = sum(text.count(kw) for kw in keywords_normalized)
    
    # カウントが1以上の場合のみ検索結果に表示する
    return (count, any(text.find(kw) != -1 for kw in keywords_normalized))

def search_doc(file_path, keywords, exclude_keywords):
    doc = win32com.client.Dispatch("Word.Application")
    doc.Visible = False
    doc.DisplayAlerts = False

    doc_open = doc.Documents.Open(file_path)
    text = doc_open.Content.Text

    doc_open.Close()
    doc.Quit()

    text = normalize_text(text.lower())  # ファイルの中身にある半角カナを全角カナに変換
    keywords_normalized = [normalize_text(kw) for kw in keywords]
    exclude_keywords_normalized = [normalize_text(kw) for kw in exclude_keywords]

    # 除外キーワードがある場合はカウントを-99999にする
    count = 0
    if any(exclude_keywords_normalized):
        for exclude_keyword in exclude_keywords_normalized:
            if exclude_keyword and normalize_text(exclude_keyword.lower()) in text:
                return (count - 99999, False)

    count = sum(text.count(kw) for kw in keywords_normalized)

    # カウントが1以上の場合のみ検索結果に表示する
    return (count, any(text.find(kw) != -1 for kw in keywords_normalized))

def search_xlsx(file_path, keywords, exclude_keywords):
    wb = load_workbook(file_path)
    text = ''
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        if isinstance(ws, openpyxl.worksheet.worksheet.Worksheet):
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        text += str(cell)
    text = normalize_text(text.lower())  # ファイルの中身にある半角カナを全角カナに変換
    keywords_normalized = [normalize_text(kw) for kw in keywords]
    exclude_keywords_normalized = [normalize_text(kw) for kw in exclude_keywords]

    # 除外キーワードがある場合はカウントを-99999にする
    count = 0
    if any(exclude_keywords_normalized):
        for exclude_keyword in exclude_keywords_normalized:
            if exclude_keyword and normalize_text(exclude_keyword.lower()) in text:
                return (count - 99999, False)

    count = sum(text.count(kw) for kw in keywords_normalized)
    
    # カウントが1以上の場合のみ検索結果に表示する
    return (count, any(text.find(kw) != -1 for kw in keywords_normalized))


def search_pdf(file_path, keywords, exclude_keywords):
    text = ''
    with open(file_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
    text = normalize_text(text.lower())  # ファイルの中身にある半角カナを全角カナに変換
    keywords_normalized = [normalize_text(kw) for kw in keywords]
    exclude_keywords_normalized = [normalize_text(kw) for kw in exclude_keywords]

    # 除外キーワードがある場合はカウントを-99999にする
    count = 0
    if any(exclude_keywords_normalized):
        for exclude_keyword in exclude_keywords_normalized:
            if exclude_keyword and normalize_text(exclude_keyword.lower()) in text:
                return (count - 99999, False)

    count = sum(text.count(kw) for kw in keywords_normalized)
    
    # カウントが1以上の場合のみ検索結果に表示する
    return (count, any(text.find(kw) != -1 for kw in keywords_normalized))

def get_selected_extensions(values):
    extensions = ['PPTX', 'DOCX', 'DOC', 'XLSX', 'PDF', 'XLS']
    selected_extensions = [ext.lower() for ext in extensions if values[f'_{ext}_']]
    
    if values['_ALL_']:
        return None
    
    return selected_extensions


def open_file(file_path):
    try:
        os.startfile(file_path)
    except OSError as e:
        print(f"Error opening file: {e}")

def main():
    sg.theme('DarkGrey8')
    layout = [
        [sg.Text('フォルダ選択')],
        [sg.InputText(key='-FOLDER_PATH-', size=(40, 1)), sg.FolderBrowse('フォルダ検索', key='-BROWSE-')],
        [sg.Text('キーワード（カンマ区切り）')],
        [sg.InputText(key='-KEYWORD-', size=(40, 1))],
        [sg.Text('除外キーワード（カンマ区切り）')],
        [sg.InputText(key='-EXCLUDE_KEYWORD-', size=(40, 1))],
        [sg.Text('拡張子の選択')],
        [sg.Checkbox('.pptx', key='_PPTX_', default=True),
        sg.Checkbox('.docx', key='_DOCX_'),
        sg.Checkbox('.doc', key='_DOC_'),
        sg.Checkbox('.xlsx', key='_XLSX_'),
        sg.Checkbox('.xls', key='_XLS_'),
        sg.Checkbox('.pdf', key='_PDF_'),
        sg.Checkbox('All', key='_ALL_')],
        [sg.Button('検索', key='-SEARCH-', size=(30, 1))],
        [sg.Text('検索結果')],
        [sg.Listbox(values=[], size=(70, 20), key='-RESULT-', enable_events=True)],
        [sg.Button('ファイルを開く', key='-OPEN-', size=(15, 1))],
    ]
    window = sg.Window('ファイル検索システム', layout)
    while True:
        event, values = window.read()
        if event == sg.WIN_CLOSED:
            break
        
        elif event == '-BROWSE-':
            folder_selected = sg.popup_get_folder('フォルダを選択してください')
            window['-FOLDER_PATH-'].update(folder_selected)
        elif event == '-SEARCH-':
            folder = values['-FOLDER_PATH-']
            keywords = values['-KEYWORD-'].split(',')
            exclude_keywords = values['-EXCLUDE_KEYWORD-'].split(',')
            if not exclude_keywords or exclude_keywords == ['']:
                exclude_keywords = []
            else:
                exclude_keywords = [normalize_text(kw) for kw in exclude_keywords]

            # チェックボックスの選択を取得
            selected_extensions = get_selected_extensions(values)

            found_files = search_files(folder, keywords, exclude_keywords, selected_extensions)
            result_list = []
            if found_files:
                for file_path, count in found_files:
                    if count > 0:
                        file_name = os.path.basename(file_path)
                        if not selected_extensions or file_path.lower().endswith(tuple(selected_extensions)):
                            result_list.append(f"{file_name} - 関連性: {count}")
                            file_paths[f"{file_name} - 関連性: {count}"] = file_path
            if not result_list:
                result_list = ['検索結果がありませんでした。']
            window['-RESULT-'].update(result_list)
        elif event == '-OPEN-':
            selected_files = values['-RESULT-']
            if selected_files:
                selected_file_name = selected_files[0]
                if selected_file_name in file_paths:
                    selected_file_path = file_paths[selected_file_name]
                    open_file(selected_file_path)
                else:
                    sg.popup_error(f"選択したファイルのパスが見つかりませんでした: {selected_file_name}")
    window.close()

if __name__ == "__main__":
    main()